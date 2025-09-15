from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.websockets import WebSocket, WebSocketDisconnect
from fastapi.responses import JSONResponse, StreamingResponse
import os
import tempfile
import zipfile
from typing import List, Dict, Optional, Union
import base64
from io import BytesIO
import asyncio
from concurrent.futures import ThreadPoolExecutor
import subprocess
import shutil
import json
import openai
import glob
from pathlib import Path
from enum import Enum
import queue
import threading
import time

# For PowerPoint processing
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
from pdf2image import convert_from_path

app = FastAPI(title="PowerPoint to Images API", version="1.0.0")

# OpenAI client - set your API key as environment variable OPENAI_API_KEY
openai_client = openai.OpenAI()

# Create thread pool for CPU-intensive tasks
executor = ThreadPoolExecutor(max_workers=4)

class PPTProcessor:
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()
    
    def extract_slide_images_libreoffice(self, ppt_path: str, output_dir: str) -> List[str]:
        """Extract slides using LibreOffice (preferred method for macOS)"""
        try:
            # Check if LibreOffice is installed
            libreoffice_path = shutil.which('libreoffice')
            if not libreoffice_path:
                # Try common macOS installation path
                libreoffice_path = '/Applications/LibreOffice.app/Contents/MacOS/soffice'
                if not os.path.exists(libreoffice_path):
                    raise Exception("LibreOffice not found. Please install with: brew install --cask libreoffice")
            
            # Convert PPT to PDF first, then PDF to images
            pdf_name = os.path.splitext(os.path.basename(ppt_path))[0] + '.pdf'
            pdf_path = os.path.join(output_dir, pdf_name)
            
            # Convert PPT/PPTX to PDF using LibreOffice
            cmd = [
                libreoffice_path if libreoffice_path.endswith('soffice') else 'libreoffice',
                '--headless', '--convert-to', 'pdf',
                '--outdir', output_dir, ppt_path
            ]
            
            result = subprocess.run(cmd, check=True, capture_output=True, text=True, timeout=60)
            
            # Verify PDF was created
            if not os.path.exists(pdf_path):
                raise Exception(f"PDF conversion failed. Expected: {pdf_path}")
            
            # Convert PDF pages to images using pdf2image
            pages = convert_from_path(
                pdf_path, 
                dpi=300,  # High DPI for quality
                fmt='PNG',
                thread_count=4,  # Use multiple threads on macOS
                poppler_path=None  # Let it auto-detect homebrew poppler
            )
            
            image_paths = []
            for i, page in enumerate(pages):
                image_path = os.path.join(output_dir, f"slide_{i+1}.png")
                # Save with high quality
                page.save(image_path, 'PNG', optimize=True, compress_level=1)
                image_paths.append(image_path)
            
            # Clean up PDF
            if os.path.exists(pdf_path):
                os.unlink(pdf_path)
            
            return image_paths
            
        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice conversion timed out")
        except subprocess.CalledProcessError as e:
            raise Exception(f"LibreOffice conversion failed: {e.stderr}")
        except Exception as e:
            raise Exception(f"LibreOffice conversion failed: {str(e)}")
    
    def extract_slide_images_manual(self, ppt_path: str, output_dir: str) -> List[str]:
        """Manual extraction using python-pptx (fallback method)"""
        try:
            prs = Presentation(ppt_path)
            image_paths = []
            
            for i, slide in enumerate(prs.slides):
                # Create a blank image for the slide
                img_width, img_height = 1920, 1080  # High quality dimensions
                img = Image.new('RGB', (img_width, img_height), 'white')
                draw = ImageDraw.Draw(img)
                
                # Process shapes in the slide
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # Extract embedded images
                        try:
                            image_stream = BytesIO(shape.image.blob)
                            shape_img = Image.open(image_stream)
                            
                            # Calculate position and size
                            left = int(shape.left.inches * (img_width / 10))  # Rough scaling
                            top = int(shape.top.inches * (img_height / 7.5))
                            width = int(shape.width.inches * (img_width / 10))
                            height = int(shape.height.inches * (img_height / 7.5))
                            
                            shape_img = shape_img.resize((width, height), Image.Resampling.LANCZOS)
                            img.paste(shape_img, (left, top))
                        except:
                            pass
                    
                    elif hasattr(shape, 'text'):
                        # Add text (basic implementation)
                        try:
                            text = shape.text
                            if text.strip():
                                left = int(shape.left.inches * (img_width / 10))
                                top = int(shape.top.inches * (img_height / 7.5))
                                
                                # Use default font for macOS
                                try:
                                    # Try common macOS system fonts
                                    font_paths = [
                                        "/System/Library/Fonts/Helvetica.ttc",
                                        "/Library/Fonts/Arial.ttf",
                                        "/System/Library/Fonts/Geneva.ttf"
                                    ]
                                    font = None
                                    for font_path in font_paths:
                                        if os.path.exists(font_path):
                                            font = ImageFont.truetype(font_path, 24)
                                            break
                                    if not font:
                                        font = ImageFont.load_default()
                                except:
                                    font = ImageFont.load_default()
                                
                                draw.text((left, top), text, fill='black', font=font)
                        except:
                            pass
                
                # Save the slide image
                image_path = os.path.join(output_dir, f"slide_{i+1}.png")
                img.save(image_path, 'PNG', quality=95, optimize=True)
                image_paths.append(image_path)
            
            return image_paths
            
        except Exception as e:
            raise Exception(f"Manual extraction failed: {str(e)}")
    
    async def process_ppt_file(self, file_path: str) -> List[str]:
        """Process PPT file and return list of image paths"""
        output_dir = tempfile.mkdtemp()
        
        # Try LibreOffice method first (best quality for macOS)
        try:
            return await asyncio.get_event_loop().run_in_executor(
                executor, self.extract_slide_images_libreoffice, file_path, output_dir
            )
        except Exception as libreoffice_error:
            print(f"LibreOffice method failed: {libreoffice_error}")
            
            # Fallback to manual method
            try:
                return await asyncio.get_event_loop().run_in_executor(
                    executor, self.extract_slide_images_manual, file_path, output_dir
                )
            except Exception as manual_error:
                raise HTTPException(
                    status_code=500, 
                    detail=f"All conversion methods failed. LibreOffice: {libreoffice_error}. Manual: {manual_error}"
                )

class ScriptGenerator:
    def __init__(self):
        self.system_prompt = """You are an expert presentation script writer and AI presenter coach. 
        Analyze the provided slide image and generate a natural, engaging presentation script.
        
        For each slide, provide:
        1. A clear, conversational script that an AI presenter would speak
        2. Key talking points and transitions
        3. Timing estimates
        4. Visual cues and emphasis points
        
        Make the script flow naturally from slide to slide, maintaining engagement and clarity.
        Keep the tone professional but approachable, suitable for an AI presenter."""
    
    async def analyze_image_with_gpt(self, image_path: str, slide_number: int, total_slides: int, previous_context: str = "") -> Dict:
        """Analyze a single slide image with GPT-4 Vision"""
        try:
            # Read and encode image
            with open(image_path, "rb") as image_file:
                image_data = base64.b64encode(image_file.read()).decode('utf-8')
            
            # Create the prompt
            user_prompt = f"""
            This is slide {slide_number} of {total_slides} in a presentation.
            
            Previous context: {previous_context}
            
            Please analyze this slide and provide:
            1. A natural presentation script (2-3 minutes of speaking)
            2. Key points to emphasize
            3. Transition to next slide (if not the last slide)
            4. Visual elements to reference
            5. Estimated speaking time in seconds
            
            Return response as JSON with this structure:
            {{
                "slide_number": {slide_number},
                "script": "Full speaking script...",
                "key_points": ["point 1", "point 2", ...],
                "visual_cues": ["visual element 1", "visual element 2", ...],
                "transition": "How to transition to next slide...",
                "estimated_time_seconds": 120,
                "slide_title": "Inferred title of the slide"
            }}
            """
            
            # Call GPT-4 Vision
            response = await asyncio.get_event_loop().run_in_executor(
                executor,
                lambda: openai_client.chat.completions.create(
                    model="gpt-4o",  # Use GPT-4 with vision
                    messages=[
                        {"role": "system", "content": self.system_prompt},
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": user_prompt},
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/png;base64,{image_data}",
                                        "detail": "high"
                                    }
                                }
                            ]
                        }
                    ],
                    max_tokens=1000,
                    temperature=0.7
                )
            )
            
            # Parse the JSON response
            content = response.choices[0].message.content
            
            # Try to extract JSON from the response
            try:
                # Look for JSON in the response
                import re
                json_match = re.search(r'\{.*\}', content, re.DOTALL)
                if json_match:
                    slide_data = json.loads(json_match.group())
                else:
                    # Fallback if no JSON found
                    slide_data = {
                        "slide_number": slide_number,
                        "script": content,
                        "key_points": [],
                        "visual_cues": [],
                        "transition": "",
                        "estimated_time_seconds": 120,
                        "slide_title": f"Slide {slide_number}"
                    }
            except json.JSONDecodeError:
                # Fallback structure
                slide_data = {
                    "slide_number": slide_number,
                    "script": content,
                    "key_points": [],
                    "visual_cues": [],
                    "transition": "",
                    "estimated_time_seconds": 120,
                    "slide_title": f"Slide {slide_number}"
                }
            
            return slide_data
            
        except Exception as e:
            print(f"Error analyzing slide {slide_number}: {e}")
            return {
                "slide_number": slide_number,
                "script": f"[Error analyzing slide {slide_number}: {str(e)}]",
                "key_points": [],
                "visual_cues": [],
                "transition": "",
                "estimated_time_seconds": 60,
                "slide_title": f"Slide {slide_number}",
                "error": str(e)
            }
    
    async def generate_presentation_script(self, image_folder: str, presentation_title: str = "Presentation") -> Dict:
        """Generate complete presentation script from folder of images"""
        
        # Get all image files from folder
        image_extensions = ['*.png', '*.jpg', '*.jpeg', '*.gif', '*.bmp', '*.webp']
        image_files = []
        
        for ext in image_extensions:
            image_files.extend(glob.glob(os.path.join(image_folder, ext)))
            image_files.extend(glob.glob(os.path.join(image_folder, ext.upper())))
        
        # Sort files naturally (slide_001.png, slide_002.png, etc.)
        image_files.sort()
        
        if not image_files:
            raise HTTPException(status_code=400, detail=f"No image files found in {image_folder}")
        
        total_slides = len(image_files)
        slides_data = []
        previous_context = ""
        
        print(f"Processing {total_slides} slides for script generation...")
        
        # Process each slide
        for i, image_path in enumerate(image_files):
            print(f"Analyzing slide {i+1}/{total_slides}: {os.path.basename(image_path)}")
            
            slide_data = await self.analyze_image_with_gpt(
                image_path, 
                i + 1, 
                total_slides, 
                previous_context
            )
            
            slides_data.append(slide_data)
            
            # Update context for next slide
            previous_context = f"Previous slide was about: {slide_data.get('slide_title', '')}. {slide_data.get('script', '')[:200]}..."
            
            # Add small delay to avoid rate limiting
            await asyncio.sleep(1)
        
        # Calculate total presentation time
        total_time = sum(slide.get('estimated_time_seconds', 120) for slide in slides_data)
        
        # Generate presentation metadata
        presentation_data = {
            "presentation_info": {
                "title": presentation_title,
                "total_slides": total_slides,
                "estimated_total_time_minutes": round(total_time / 60, 1),
                "estimated_total_time_seconds": total_time,
                "generated_at": asyncio.get_event_loop().time(),
                "source_folder": image_folder
            },
            "slides": slides_data,
            "presentation_notes": {
                "introduction": f"Welcome to {presentation_title}. This presentation contains {total_slides} slides and is estimated to take {round(total_time / 60, 1)} minutes.",
                "conclusion": "Thank you for your attention. Are there any questions?",
                "ai_presenter_tips": [
                    "Maintain steady pacing throughout",
                    "Emphasize key points with vocal inflection",
                    "Pause briefly between major sections",
                    "Reference visual elements when mentioned"
                ]
            }
        }
        
        return presentation_data

class PresentationState(Enum):
    IDLE = "idle"
    PRESENTING = "presenting"
    INTERRUPTED = "interrupted"
    PAUSED = "paused"

class TTSChunk:
    def __init__(self, chunk_id: int, audio_data: bytes, is_final: bool = False):
        self.chunk_id = chunk_id
        self.audio_data = audio_data
        self.is_final = is_final

class PresentationManager:
    def __init__(self):
        self.presentation_data: Optional[Dict] = None
        self.current_slide: int = 0
        self.state: PresentationState = PresentationState.IDLE
        self.audio_queue: queue.Queue = queue.Queue()
        self.current_chunk_id: int = 0
        self.paused_chunk_id: Optional[int] = None
        self.websocket: Optional[WebSocket] = None
        self.tts_task: Optional[asyncio.Task] = None
        
        # TTS Configuration
        self.tts_voice = "alloy"  # Options: alloy, echo, fable, onyx, nova, shimmer
        self.tts_model = "tts-1"  # Options: tts-1, tts-1-hd (higher quality)
        self.tts_speed = 1.0      # Speed: 0.25 to 4.0
        
    def configure_tts(self, voice: str = "alloy", model: str = "tts-1", speed: float = 1.0):
        """Configure TTS settings"""
        valid_voices = ["alloy", "echo", "fable", "onyx", "nova", "shimmer"]
        valid_models = ["tts-1", "tts-1-hd"]
        
        if voice not in valid_voices:
            raise ValueError(f"Invalid voice. Choose from: {valid_voices}")
        if model not in valid_models:
            raise ValueError(f"Invalid model. Choose from: {valid_models}")
        if not (0.25 <= speed <= 4.0):
            raise ValueError("Speed must be between 0.25 and 4.0")
            
        self.tts_voice = voice
        self.tts_model = model
        self.tts_speed = speed
        
    def load_presentation(self, presentation_data: Dict):
        """Load presentation data"""
        self.presentation_data = presentation_data
        self.current_slide = 0
        self.state = PresentationState.IDLE
        
    async def start_slide(self, slide_number: int, slide_image: str = None):
        """Start presenting a specific slide with debugging"""
        print(f"Starting slide {slide_number}")
        
        if not self.presentation_data:
            raise ValueError("No presentation data loaded")
            
        if slide_number > len(self.presentation_data.get('slides', [])):
            raise ValueError(f"Slide {slide_number} does not exist")
            
        self.current_slide = slide_number
        self.state = PresentationState.PRESENTING
        
        # Get slide script
        slide_data = self.presentation_data['slides'][slide_number - 1]
        script = slide_data.get('script', '')
        print(f"Slide script length: {len(script)} characters")
        
        # Generate TTS chunks for the script
        await self._generate_tts_chunks(script, slide_number)
        
        # Start sending audio chunks
        print("Creating TTS task")
        self.tts_task = asyncio.create_task(self._send_audio_chunks())
        print(f"TTS task created: {self.tts_task}")
        
    async def _generate_tts_chunks(self, text: str, slide_number: int):
        """Generate TTS audio chunks from text using OpenAI TTS with debugging"""
        print(f"Starting TTS generation for slide {slide_number}")
        
        # Clear the queue
        while not self.audio_queue.empty():
            try:
                self.audio_queue.get_nowait()
            except queue.Empty:
                break
                
        self.current_chunk_id = 0
        
        # Split text into manageable chunks for TTS
        chunks = self._split_text_for_tts(text)
        print(f"Split text into {len(chunks)} chunks")
        
        for i, chunk_text in enumerate(chunks):
            if chunk_text.strip():
                try:
                    print(f"Generating TTS for chunk {i+1}/{len(chunks)}")
                    audio_data = await self._text_to_speech(chunk_text.strip())
                    print(f"Generated {len(audio_data)} bytes of audio")
                    
                    is_final = (i == len(chunks) - 1)
                    chunk = TTSChunk(self.current_chunk_id, audio_data, is_final)
                    self.audio_queue.put(chunk)
                    print(f"Added chunk {self.current_chunk_id} to queue (final: {is_final})")
                    self.current_chunk_id += 1
                    
                except Exception as e:
                    print(f"Error generating TTS for chunk {i}: {e}")
    
    def _split_text_for_tts(self, text: str, max_chars: int = 500) -> List[str]:
        """Split text into small chunks at punctuation marks for smaller TTS files"""
        import re
        
        # Split at sentence-ending punctuation: . ! ? and also : ;
        sentences = re.split(r'[.!?:;]+', text)
        
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            # If adding this sentence would exceed limit, save current chunk
            if len(current_chunk) + len(sentence) + 2 > max_chars and current_chunk:
                chunks.append(current_chunk.strip() + ".")
                current_chunk = sentence
            else:
                if current_chunk:
                    current_chunk += ". " + sentence
                else:
                    current_chunk = sentence
        
        # Add the last chunk
        if current_chunk:
            chunks.append(current_chunk.strip() + ".")
        
        # If chunks are still too long, split further at commas
        final_chunks = []
        for chunk in chunks:
            if len(chunk) > max_chars:
                # Split at commas
                parts = chunk.split(',')
                sub_chunk = ""
                
                for part in parts:
                    part = part.strip()
                    if len(sub_chunk) + len(part) + 2 > max_chars and sub_chunk:
                        final_chunks.append(sub_chunk.strip() + ".")
                        sub_chunk = part
                    else:
                        if sub_chunk:
                            sub_chunk += ", " + part
                        else:
                            sub_chunk = part
                
                if sub_chunk:
                    final_chunks.append(sub_chunk.strip() + ".")
            else:
                final_chunks.append(chunk)
        
        return final_chunks
                
    async def _text_to_speech(self, text: str) -> bytes:
        """Convert text to speech using OpenAI TTS API"""
        try:
            # Call OpenAI TTS API
            response = await asyncio.get_event_loop().run_in_executor(
                executor,
                lambda: openai_client.audio.speech.create(
                    model=self.tts_model,
                    voice=self.tts_voice,
                    input=text,
                    speed=self.tts_speed,
                    response_format="mp3"  # Options: mp3, opus, aac, flac
                )
            )
            
            # Get audio content as bytes
            audio_bytes = response.content
            return audio_bytes
            
        except Exception as e:
            print(f"OpenAI TTS API error: {e}")
            
            # Fallback: return error message as text bytes
            error_message = f"[TTS Error: {str(e)}]"
            return error_message.encode('utf-8')
        
    async def _send_audio_chunks(self):
        """Send audio chunks through WebSocket with size limits"""
        try:
            while self.state == PresentationState.PRESENTING:
                if not self.audio_queue.empty():
                    chunk = self.audio_queue.get()
                    
                    # Encode audio to base64 first to check size
                    audio_b64 = base64.b64encode(chunk.audio_data).decode('utf-8')
                    max_size = 500000

                    if len(audio_b64) > max_size:
                        # Split large audio into smaller parts
                        num_parts = (len(audio_b64) + max_size - 1) // max_size
                        print(f"Splitting large chunk {chunk.chunk_id} into {num_parts} parts")
                        
                        for part_idx in range(num_parts):
                            start_idx = part_idx * max_size
                            end_idx = min((part_idx + 1) * max_size, len(audio_b64))
                            part_data = audio_b64[start_idx:end_idx]
                            
                            is_final_part = (part_idx == num_parts - 1) and chunk.is_final
                            
                            if self.websocket:
                                await self.websocket.send_json({
                                    "type": "audio_chunk",
                                    "chunk_id": f"{chunk.chunk_id}-{part_idx}",
                                    "audio_data": part_data,
                                    "slide_number": self.current_slide,
                                    "is_final": is_final_part,
                                    "part": part_idx + 1,
                                    "total_parts": num_parts
                                })
                                print(f"Sent part {part_idx + 1}/{num_parts} of chunk {chunk.chunk_id}")
                            
                            await asyncio.sleep(0.2)  # Small delay between parts
                    else:
                        # Send normal sized chunk
                        if self.websocket:
                            await self.websocket.send_json({
                                "type": "audio_chunk",
                                "chunk_id": chunk.chunk_id,
                                "audio_data": audio_b64,
                                "slide_number": self.current_slide,
                                "is_final": chunk.is_final
                            })
                            print(f"Sent chunk {chunk.chunk_id} ({len(audio_b64)} chars)")
                    
                    if chunk.is_final:
                        await self._send_slide_done()
                        break
                        
                    await asyncio.sleep(0.5)
                else:
                    await asyncio.sleep(0.1)
                    
        except Exception as e:
            print(f"Error sending audio chunks: {e}")
            import traceback
            traceback.print_exc()
            
    async def _send_slide_done(self):
        """Send slide done signal"""
        if self.websocket:
            await self.websocket.send_json({
                "type": "slide_done",
                "slide_number": self.current_slide,
                "message": f"Slide {self.current_slide} completed"
            })
        self.state = PresentationState.IDLE
        
    async def handle_interrupt(self, question: str = ""):
        """Handle user interrupt for Q&A"""
        if self.state != PresentationState.PRESENTING:
            return
            
        # Pause the current presentation
        self.state = PresentationState.INTERRUPTED
        
        # Cancel current TTS task
        if self.tts_task and not self.tts_task.done():
            self.tts_task.cancel()
            
        # Save current position
        self.paused_chunk_id = self.current_chunk_id
        
        # Send interrupt acknowledgment
        if self.websocket:
            await self.websocket.send_json({
                "type": "interrupt_acknowledged",
                "message": "Presentation paused for Q&A",
                "paused_at_chunk": self.paused_chunk_id
            })
            
        # Handle Q&A (placeholder function)
        answer = await self._qa_function(question)
        
        # Send Q&A response as TTS
        qa_audio = await self._text_to_speech(answer)
        if self.websocket:
            await self.websocket.send_json({
                "type": "qa_response",
                "question": question,
                "answer": answer,
                "audio_data": base64.b64encode(qa_audio).decode('utf-8')
            })
            
        # Resume presentation
        await self._resume_presentation()
        
    async def _qa_function(self, question: str) -> str:
        """Handle Q&A using GPT with context from the current slide"""
        try:
            if not self.presentation_data or not question.strip():
                return "Thank you for your question. Let me continue with the presentation and we can discuss this further at the end."
            
            # Get current slide context
            current_slide_data = self.presentation_data['slides'][self.current_slide - 1] if self.current_slide > 0 else {}
            slide_context = current_slide_data.get('script', '')
            slide_title = current_slide_data.get('slide_title', f'Slide {self.current_slide}')
            
            # Create context-aware prompt
            system_prompt = """You are an AI presenter assistant. A user has interrupted the presentation with a question. 
            Provide a helpful, concise answer based on the current slide content and general knowledge. 
            Keep the response brief (1-2 sentences) as this is an interruption during a live presentation.
            If the question is not directly related to the current slide, acknowledge it and suggest discussing it later."""
            
            user_prompt = f"""
            Current slide: "{slide_title}"
            Slide content: "{slide_context[:500]}..."
            
            User question: "{question}"
            
            Please provide a brief, helpful response.
            """
            
            # Call OpenAI API
            response = await asyncio.get_event_loop().run_in_executor(
                executor,
                lambda: openai_client.chat.completions.create(
                    model="gpt-4o-mini",  # Use faster model for real-time responses
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ],
                    max_tokens=150,  # Keep responses concise
                    temperature=0.7
                )
            )
            
            answer = response.choices[0].message.content
            return answer
            
        except Exception as e:
            print(f"Q&A function error: {e}")
            return "Thank you for your question. Let me continue with the presentation and we can address this at the end."
        
    async def _resume_presentation(self):
        """Resume presentation from where it was paused"""
        if self.websocket:
            await self.websocket.send_json({
                "type": "presentation_resumed",
                "message": "Resuming presentation",
                "resuming_from_chunk": self.paused_chunk_id
            })
            
        # Resume sending audio chunks
        self.state = PresentationState.PRESENTING
        self.tts_task = asyncio.create_task(self._send_audio_chunks())
        
    def stop_presentation(self):
        """Stop the current presentation"""
        self.state = PresentationState.IDLE
        if self.tts_task and not self.tts_task.done():
            self.tts_task.cancel()
            
        # Clear the queue
        while not self.audio_queue.empty():
            try:
                self.audio_queue.get_nowait()
            except queue.Empty:
                break

# Initialize processors
processor = PPTProcessor()
script_generator = ScriptGenerator()
presentation_manager = PresentationManager()

@app.post("/convert-ppt/")
async def convert_ppt_to_images(file: UploadFile = File(...)):
    """
    Convert PowerPoint file to high-quality images
    Returns JSON with base64 encoded images for each slide
    """
    
    # Validate file type
    if not file.filename.lower().endswith(('.ppt', '.pptx')):
        raise HTTPException(status_code=400, detail="File must be a PowerPoint file (.ppt or .pptx)")
    
    # Create temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        content = await file.read()
        tmp_file.write(content)
        tmp_file_path = tmp_file.name
    
    try:
        # Process the PowerPoint file
        image_paths = await processor.process_ppt_file(tmp_file_path)
        
        # Convert images to base64
        slides_data = []
        for i, image_path in enumerate(image_paths):
            with open(image_path, 'rb') as img_file:
                img_data = img_file.read()
                img_base64 = base64.b64encode(img_data).decode('utf-8')
                
                slides_data.append({
                    "slide_number": i + 1,
                    "image_data": img_base64,
                    "content_type": "image/png"
                })
            
            # Clean up individual image file
            os.unlink(image_path)
        
        return JSONResponse(content={
            "status": "success",
            "total_slides": len(slides_data),
            "slides": slides_data
        })
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    
    finally:
        # Clean up temporary file
        if os.path.exists(tmp_file_path):
            os.unlink(tmp_file_path)

@app.post("/generate-script/")
async def generate_presentation_script(
    images_folder: str = Form(...),
    presentation_title: str = Form("AI Generated Presentation"),
    openai_api_key: Optional[str] = Form(None)
):
    """
    Generate presentation script from folder of slide images using GPT-4 Vision
    
    Args:
        images_folder: Path to folder containing slide images
        presentation_title: Title of the presentation
        openai_api_key: OpenAI API key (optional if set as environment variable)
    
    Returns:
        JSON with complete presentation script and metadata
    """
    
    # Set OpenAI API key if provided
    if openai_api_key:
        openai_client.api_key = openai_api_key
    
    # Validate folder exists
    if not os.path.exists(images_folder):
        raise HTTPException(status_code=400, detail=f"Folder not found: {images_folder}")
    
    if not os.path.isdir(images_folder):
        raise HTTPException(status_code=400, detail=f"Path is not a directory: {images_folder}")
    
    try:
        # Generate the presentation script
        presentation_data = await script_generator.generate_presentation_script(
            images_folder, 
            presentation_title
        )
        
        return JSONResponse(content={
            "status": "success",
            "message": "Presentation script generated successfully",
            **presentation_data
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Script generation failed: {str(e)}")

@app.post("/upload-slides-generate-script/")
async def upload_slides_and_generate_script(
    file: UploadFile = File(...),
    presentation_title: str = Form("AI Generated Presentation"),
    openai_api_key: Optional[str] = Form(None)
):
    """
    Upload convert-ppt response JSON and generate presentation script
    
    Args:
        file: JSON file containing convert-ppt response with slides array
        presentation_title: Title of the presentation
        openai_api_key: OpenAI API key (optional if set as environment variable)
    
    Returns:
        JSON with complete presentation script and metadata
    """
    
    # Set OpenAI API key if provided
    if openai_api_key:
        openai_client.api_key = openai_api_key
    
    # Validate file type
    if not file.filename.lower().endswith('.json'):
        raise HTTPException(status_code=400, detail="File must be a JSON file")
    
    try:
        # Read and parse JSON
        content = await file.read()
        convert_ppt_response = json.loads(content.decode('utf-8'))
        
        # Validate structure
        if 'slides' not in convert_ppt_response:
            raise HTTPException(status_code=400, detail="Invalid format: missing 'slides' key")
        
        slides = convert_ppt_response.get('slides', [])
        if not slides:
            raise HTTPException(status_code=400, detail="No slides found in response")
        
        # Create temporary directory for images
        temp_dir = tempfile.mkdtemp()
        
        try:
            # Save base64 images to files
            image_files = []
            for slide in slides:
                slide_number = slide.get('slide_number', len(image_files) + 1)
                image_data = slide.get('image_data')
                
                if not image_data:
                    continue
                
                # Decode base64 image
                image_bytes = base64.b64decode(image_data)
                
                # Save to file
                filename = f"slide_{slide_number:03d}.png"
                filepath = os.path.join(temp_dir, filename)
                
                with open(filepath, 'wb') as f:
                    f.write(image_bytes)
                image_files.append(filepath)
            
            if not image_files:
                raise HTTPException(status_code=400, detail="No valid image data found in slides")
            
            print(f"Saved {len(image_files)} images to {temp_dir}")
            
            # Generate presentation script
            presentation_data = await script_generator.generate_presentation_script(
                temp_dir,
                presentation_title
            )
            
            return JSONResponse(content={
                "status": "success",
                "message": "Slides processed and presentation script generated successfully",
                **presentation_data
            })
            
        finally:
            # Clean up temporary files
            try:
                shutil.rmtree(temp_dir)
            except:
                pass
        
    except json.JSONDecodeError:
        raise HTTPException(status_code=400, detail="Invalid JSON format")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Script generation failed: {str(e)}")

@app.post("/test-tts/")
async def test_tts(text: str = Form("This is a test of the text to speech system.")):
    """
    Test TTS generation with custom text and return WAV file
    
    Args:
        text: Text to convert to speech
    
    Returns:
        WAV audio file
    """
    try:
        # Generate TTS audio using OpenAI
        response = openai_client.audio.speech.create(
            model="tts-1",
            voice="alloy", 
            input=text,
            response_format="wav"  # Return WAV instead of MP3
        )
        
        # Get audio content as bytes
        audio_bytes = response.content
        
        # Return as streaming response with proper headers
        return StreamingResponse(
            BytesIO(audio_bytes),
            media_type="audio/wav",
            headers={
                "Content-Disposition": f"attachment; filename=tts_output.wav",
                "Content-Length": str(len(audio_bytes))
            }
        )
        
    except Exception as e:
        raise HTTPException(
            status_code=500, 
            detail=f"TTS generation failed: {str(e)}"
        )

@app.websocket("/ws/presentation")
async def presentation_websocket(websocket: WebSocket):
    """
    WebSocket endpoint for real-time presentation with TTS audio streaming
    
    Protocol:
    Client -> Server:
    - {"type": "load_presentation", "data": presentation_json}
    - {"type": "slide_start", "slide_number": 1, "slide_image": "base64_image"}
    - {"type": "interrupt", "question": "optional question text"}
    - {"type": "configure_tts", "voice": "alloy", "model": "tts-1", "speed": 1.0}
    - {"type": "resume"}
    - {"type": "stop"}
    - {"type": "get_status"}
    - {"type": "ping"}
    
    Server -> Client:
    - {"type": "connected"}
    - {"type": "presentation_loaded"}
    - {"type": "slide_started"}
    - {"type": "audio_chunk", "chunk_id": 1, "audio_data": "base64_audio", "slide_number": 1}
    - {"type": "slide_done", "slide_number": 1}
    - {"type": "interrupt_acknowledged"}
    - {"type": "qa_response", "question": "...", "answer": "...", "audio_data": "base64_audio"}
    - {"type": "presentation_resumed"}
    - {"type": "tts_configured"}
    - {"type": "status"}
    - {"type": "pong"}
    - {"type": "error", "message": "error description"}
    """
    
    print(f"WebSocket connection attempt from {websocket.client}")
    
    try:
        await websocket.accept()
        print("WebSocket connection accepted")
        
        # Set websocket in presentation manager
        presentation_manager.websocket = websocket
        
        # Send initial connection confirmation
        await websocket.send_json({
            "type": "connected",
            "message": "WebSocket connected successfully",
            "server_time": time.time(),
            "protocol_version": "1.0"
        })
        print("Sent connection confirmation")
        
        # Main message handling loop
        while True:
            try:
                # Wait for message with timeout
                data = await asyncio.wait_for(websocket.receive_json(), timeout=60.0)
                message_type = data.get("type", "unknown")
                
                print(f"Received message: {message_type}")
                
                # Handle different message types
                if message_type == "ping":
                    await websocket.send_json({
                        "type": "pong",
                        "timestamp": time.time()
                    })
                    print("Sent pong response")
                
                elif message_type == "load_presentation":
                    try:
                        presentation_data = data.get("data")
                        if not presentation_data:
                            await websocket.send_json({
                                "type": "error",
                                "message": "No presentation data provided"
                            })
                            continue
                            
                        if "slides" not in presentation_data:
                            await websocket.send_json({
                                "type": "error",
                                "message": "Invalid presentation format: missing 'slides' key"
                            })
                            continue
                        
                        presentation_manager.load_presentation(presentation_data)
                        
                        await websocket.send_json({
                            "type": "presentation_loaded",
                            "message": f"Loaded presentation with {len(presentation_data.get('slides', []))} slides",
                            "total_slides": len(presentation_data.get('slides', [])),
                            "title": presentation_data.get('presentation_info', {}).get('title', 'Untitled')
                        })
                        print(f"Loaded presentation: {len(presentation_data.get('slides', []))} slides")
                        
                    except Exception as e:
                        print(f"Error loading presentation: {e}")
                        await websocket.send_json({
                            "type": "error",
                            "message": f"Failed to load presentation: {str(e)}"
                        })
                
                elif message_type == "slide_start":
                    try:
                        slide_number = data.get("slide_number")
                        slide_image = data.get("slide_image")  # Optional base64 image
                        
                        if not slide_number:
                            await websocket.send_json({
                                "type": "error",
                                "message": "Slide number is required"
                            })
                            continue
                        
                        if not presentation_manager.presentation_data:
                            await websocket.send_json({
                                "type": "error",
                                "message": "No presentation loaded. Load presentation first."
                            })
                            continue
                        
                        total_slides = len(presentation_manager.presentation_data.get('slides', []))
                        if slide_number > total_slides or slide_number < 1:
                            await websocket.send_json({
                                "type": "error",
                                "message": f"Invalid slide number. Must be between 1 and {total_slides}"
                            })
                            continue
                        
                        print(f"Starting slide {slide_number}")
                        
                        # Send immediate acknowledgment
                        await websocket.send_json({
                            "type": "slide_started",
                            "slide_number": slide_number,
                            "message": f"Started slide {slide_number}",
                            "timestamp": time.time()
                        })
                        
                        # Start the slide presentation
                        await presentation_manager.start_slide(slide_number, slide_image)
                        
                    except ValueError as e:
                        print(f"Slide start validation error: {e}")
                        await websocket.send_json({
                            "type": "error",
                            "message": str(e)
                        })
                    except Exception as e:
                        print(f"Error starting slide: {e}")
                        import traceback
                        traceback.print_exc()
                        await websocket.send_json({
                            "type": "error",
                            "message": f"Failed to start slide: {str(e)}"
                        })
                
                elif message_type == "interrupt":
                    try:
                        question = data.get("question", "")
                        print(f"Interrupt received with question: {question[:50]}...")
                        
                        if presentation_manager.state != PresentationState.PRESENTING:
                            await websocket.send_json({
                                "type": "error",
                                "message": f"Cannot interrupt. Current state: {presentation_manager.state.value}"
                            })
                            continue
                        
                        await presentation_manager.handle_interrupt(question)
                        
                    except Exception as e:
                        print(f"Error handling interrupt: {e}")
                        await websocket.send_json({
                            "type": "error",
                            "message": f"Failed to handle interrupt: {str(e)}"
                        })
                
                elif message_type == "configure_tts":
                    try:
                        voice = data.get("voice", "alloy")
                        model = data.get("model", "tts-1")
                        speed = data.get("speed", 1.0)
                        
                        presentation_manager.configure_tts(voice, model, speed)
                        
                        await websocket.send_json({
                            "type": "tts_configured",
                            "message": f"TTS configured: voice={voice}, model={model}, speed={speed}",
                            "voice": voice,
                            "model": model,
                            "speed": speed
                        })
                        print(f"TTS configured: {voice}, {model}, {speed}")
                        
                    except ValueError as e:
                        await websocket.send_json({
                            "type": "error",
                            "message": f"TTS configuration error: {str(e)}"
                        })
                    except Exception as e:
                        print(f"Error configuring TTS: {e}")
                        await websocket.send_json({
                            "type": "error",
                            "message": f"Failed to configure TTS: {str(e)}"
                        })
                
                elif message_type == "resume":
                    try:
                        if presentation_manager.state == PresentationState.INTERRUPTED:
                            await presentation_manager._resume_presentation()
                        else:
                            await websocket.send_json({
                                "type": "error",
                                "message": f"Cannot resume. Current state: {presentation_manager.state.value}"
                            })
                    except Exception as e:
                        print(f"Error resuming: {e}")
                        await websocket.send_json({
                            "type": "error",
                            "message": f"Failed to resume: {str(e)}"
                        })
                
                elif message_type == "stop":
                    try:
                        presentation_manager.stop_presentation()
                        await websocket.send_json({
                            "type": "presentation_stopped",
                            "message": "Presentation stopped",
                            "timestamp": time.time()
                        })
                        print("Presentation stopped")
                    except Exception as e:
                        print(f"Error stopping: {e}")
                        await websocket.send_json({
                            "type": "error",
                            "message": f"Failed to stop: {str(e)}"
                        })
                
                elif message_type == "get_status":
                    try:
                        total_slides = 0
                        current_slide_info = None
                        
                        if presentation_manager.presentation_data:
                            slides = presentation_manager.presentation_data.get('slides', [])
                            total_slides = len(slides)
                            
                            if presentation_manager.current_slide > 0 and presentation_manager.current_slide <= total_slides:
                                current_slide_info = slides[presentation_manager.current_slide - 1]
                        
                        await websocket.send_json({
                            "type": "status",
                            "state": presentation_manager.state.value,
                            "current_slide": presentation_manager.current_slide,
                            "total_slides": total_slides,
                            "current_slide_title": current_slide_info.get('slide_title', '') if current_slide_info else '',
                            "tts_config": {
                                "voice": presentation_manager.tts_voice,
                                "model": presentation_manager.tts_model,
                                "speed": presentation_manager.tts_speed
                            },
                            "timestamp": time.time()
                        })
                        print("Sent status update")
                    except Exception as e:
                        print(f"Error getting status: {e}")
                        await websocket.send_json({
                            "type": "error",
                            "message": f"Failed to get status: {str(e)}"
                        })
                
                else:
                    await websocket.send_json({
                        "type": "error",
                        "message": f"Unknown message type: {message_type}"
                    })
                    print(f"Unknown message type: {message_type}")
            
            except asyncio.TimeoutError:
                print("WebSocket receive timeout - sending keepalive")
                try:
                    await websocket.send_json({
                        "type": "keepalive",
                        "timestamp": time.time()
                    })
                except:
                    print("Failed to send keepalive - connection likely closed")
                    break
            
            except json.JSONDecodeError as e:
                print(f"JSON decode error: {e}")
                await websocket.send_json({
                    "type": "error",
                    "message": f"Invalid JSON: {str(e)}"
                })
            
            except Exception as e:
                print(f"Error processing message: {e}")
                import traceback
                traceback.print_exc()
                try:
                    await websocket.send_json({
                        "type": "error",
                        "message": f"Message processing error: {str(e)}"
                    })
                except:
                    print("Failed to send error response - connection likely closed")
                    break
    
    except WebSocketDisconnect:
        print("WebSocket disconnected normally")
    except Exception as e:
        print(f"WebSocket connection error: {e}")
        import traceback
        traceback.print_exc()
    finally:
        print("Cleaning up WebSocket connection")
        # Clean up
        if presentation_manager.websocket == websocket:
            presentation_manager.stop_presentation()
            presentation_manager.websocket = None
        print("WebSocket cleanup completed")

@app.post("/upload-presentation-script/")
async def upload_presentation_script(file: UploadFile = File(...)):
    """
    Upload a presentation script JSON file to be used with WebSocket presentation
    
    Args:
        file: JSON file containing presentation script (from /generate-script/ endpoint)
    
    Returns:
        Confirmation that script is ready for WebSocket presentation
    """
    
    if not file.filename.lower().endswith('.json'):
        raise HTTPException(status_code=400, detail="File must be a JSON file")
    
    try:
        # Read and parse JSON
        content = await file.read()
        presentation_data = json.loads(content.decode('utf-8'))
        
        # Validate structure
        if 'slides' not in presentation_data:
            raise HTTPException(status_code=400, detail="Invalid presentation format: missing 'slides' key")
        
        # Load into presentation manager
        presentation_manager.load_presentation(presentation_data)
        
        return JSONResponse(content={
            "status": "success",
            "message": "Presentation script uploaded and ready",
            "total_slides": len(presentation_data.get('slides', [])),
            "websocket_endpoint": "/ws/presentation",
            "usage_instructions": {
                "1": "Connect to WebSocket at /ws/presentation",
                "2": "Send slide_start messages with slide numbers",
                "3": "Use interrupt messages for Q&A",
                "4": "Server will stream TTS audio chunks"
            }
        })
        
    except json.JSONDecodeError:
        raise HTTPException(status_code=400, detail="Invalid JSON format")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing presentation script: {str(e)}")

@app.get("/")
async def root():
    return {"message": "PowerPoint to Images API (macOS)", "version": "1.0.0"}

@app.get("/health")
async def health_check():
    return {"status": "healthy", "platform": "macOS"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)